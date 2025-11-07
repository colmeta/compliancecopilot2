"""
CLARITY Funding Engine - Document Converter
Converts Markdown documents to PDF, Word, and PowerPoint formats
Presidential-grade formatting and layout
"""

import os
import io
import re
from typing import Dict, List, Optional, Tuple
from datetime import datetime

# PDF Generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image
from reportlab.pdfgen import canvas

# Word Generation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

# PowerPoint Generation
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

import markdown2
import logging

logger = logging.getLogger(__name__)


class DocumentConverter:
    """Convert Markdown documents to professional PDF, Word, and PowerPoint formats"""
    
    def __init__(self):
        self.brand_color = colors.HexColor('#F59E0B')  # Amber brand color
        self.dark_bg = colors.HexColor('#0F172A')  # Slate 900
        self.text_color = colors.HexColor('#1E293B')  # Slate 800
        
    def convert_to_pdf(self, markdown_content: str, output_path: str, metadata: Dict = None) -> str:
        """Convert Markdown to professionally formatted PDF"""
        try:
            # Create PDF
            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=1*inch,
                bottomMargin=0.75*inch
            )
            
            # Styles
            styles = getSampleStyleSheet()
            
            # Custom styles for presidential look
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=self.text_color,
                spaceAfter=30,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold'
            )
            
            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=self.brand_color,
                spaceAfter=12,
                spaceBefore=12,
                fontName='Helvetica-Bold'
            )
            
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=self.text_color,
                spaceAfter=10,
                spaceBefore=10,
                fontName='Helvetica-Bold'
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['BodyText'],
                fontSize=11,
                textColor=self.text_color,
                spaceAfter=12,
                alignment=TA_JUSTIFY,
                fontName='Helvetica'
            )
            
            # Parse markdown and build content
            story = []
            
            # Add header with metadata
            if metadata:
                header_text = f"""
                <b>{metadata.get('company_name', 'Company Name')}</b><br/>
                {metadata.get('document_type', 'Document')}<br/>
                Generated: {datetime.now().strftime('%B %d, %Y')}
                """
                story.append(Paragraph(header_text, body_style))
                story.append(Spacer(1, 0.3*inch))
            
            # Parse markdown line by line
            lines = markdown_content.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                
                if not line:
                    story.append(Spacer(1, 0.1*inch))
                    i += 1
                    continue
                
                # H1 Title
                if line.startswith('# '):
                    text = line[2:].strip()
                    story.append(Paragraph(text, title_style))
                    story.append(Spacer(1, 0.2*inch))
                
                # H2 Heading
                elif line.startswith('## '):
                    text = line[3:].strip()
                    story.append(Paragraph(text, heading1_style))
                
                # H3 Heading
                elif line.startswith('### '):
                    text = line[4:].strip()
                    story.append(Paragraph(text, heading2_style))
                
                # Bullet list
                elif line.startswith('- ') or line.startswith('* '):
                    text = '• ' + line[2:].strip()
                    story.append(Paragraph(text, body_style))
                
                # Table (simple detection)
                elif '|' in line and i + 1 < len(lines) and '|' in lines[i+1]:
                    # Parse table
                    table_lines = [line]
                    i += 1
                    while i < len(lines) and '|' in lines[i]:
                        table_lines.append(lines[i])
                        i += 1
                    
                    # Create table (basic implementation)
                    table_data = []
                    for tline in table_lines:
                        if not tline.strip().startswith('|--'):  # Skip separator lines
                            cells = [cell.strip() for cell in tline.split('|')[1:-1]]
                            table_data.append(cells)
                    
                    if table_data:
                        t = Table(table_data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), self.brand_color),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 12),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('GRID', (0, 0), (-1, -1), 1, colors.grey)
                        ]))
                        story.append(t)
                        story.append(Spacer(1, 0.2*inch))
                    continue
                
                # Regular paragraph
                else:
                    # Convert markdown bold/italic
                    text = line
                    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
                    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
                    story.append(Paragraph(text, body_style))
                
                i += 1
            
            # Build PDF
            doc.build(story)
            logger.info(f"PDF generated: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"PDF conversion failed: {e}")
            raise
    
    def convert_to_word(self, markdown_content: str, output_path: str, metadata: Dict = None) -> str:
        """Convert Markdown to professionally formatted Word document"""
        try:
            doc = Document()
            
            # Set document properties
            if metadata:
                doc.core_properties.title = metadata.get('document_type', 'Document')
                doc.core_properties.author = metadata.get('company_name', 'Company')
                doc.core_properties.created = datetime.now()
            
            # Add header
            if metadata:
                header_para = doc.add_paragraph()
                header_para.add_run(metadata.get('company_name', 'Company Name')).bold = True
                header_para.add_run(f"\n{metadata.get('document_type', 'Document')}")
                header_para.add_run(f"\nGenerated: {datetime.now().strftime('%B %d, %Y')}")
                header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                doc.add_paragraph()  # Spacer
            
            # Parse markdown
            lines = markdown_content.split('\n')
            i = 0
            
            while i < len(lines):
                line = lines[i].strip()
                
                if not line:
                    doc.add_paragraph()  # Empty line
                    i += 1
                    continue
                
                # H1 Title
                if line.startswith('# '):
                    text = line[2:].strip()
                    heading = doc.add_heading(text, level=1)
                    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                # H2 Heading
                elif line.startswith('## '):
                    text = line[3:].strip()
                    doc.add_heading(text, level=2)
                
                # H3 Heading
                elif line.startswith('### '):
                    text = line[4:].strip()
                    doc.add_heading(text, level=3)
                
                # Bullet list
                elif line.startswith('- ') or line.startswith('* '):
                    text = line[2:].strip()
                    doc.add_paragraph(text, style='List Bullet')
                
                # Table
                elif '|' in line and i + 1 < len(lines) and '|' in lines[i+1]:
                    # Parse table
                    table_lines = [line]
                    i += 1
                    while i < len(lines) and '|' in lines[i]:
                        table_lines.append(lines[i])
                        i += 1
                    
                    # Create table
                    table_data = []
                    for tline in table_lines:
                        if not tline.strip().startswith('|--'):
                            cells = [cell.strip() for cell in tline.split('|')[1:-1]]
                            table_data.append(cells)
                    
                    if table_data:
                        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                        table.style = 'Light Grid Accent 1'
                        
                        for row_idx, row_data in enumerate(table_data):
                            for col_idx, cell_data in enumerate(row_data):
                                table.rows[row_idx].cells[col_idx].text = cell_data
                    
                    continue
                
                # Regular paragraph
                else:
                    para = doc.add_paragraph()
                    # Parse bold/italic
                    text = line
                    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            para.add_run(part[2:-2]).bold = True
                        elif part.startswith('*') and part.endswith('*'):
                            para.add_run(part[1:-1]).italic = True
                        else:
                            para.add_run(part)
                
                i += 1
            
            # Save
            doc.save(output_path)
            logger.info(f"Word document generated: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Word conversion failed: {e}")
            raise
    
    def convert_to_powerpoint(self, markdown_content: str, output_path: str, metadata: Dict = None) -> str:
        """Convert Markdown to professionally formatted PowerPoint (for Pitch Deck)"""
        try:
            prs = Presentation()
            prs.slide_width = PptxInches(10)
            prs.slide_height = PptxInches(7.5)
            
            # Parse markdown into slides (each ## is a new slide)
            lines = markdown_content.split('\n')
            current_slide_title = None
            current_slide_content = []
            
            def add_slide(title, content):
                """Add a slide with title and content"""
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                title_shape = slide.shapes.title
                title_shape.text = title
                
                # Content
                content_box = slide.shapes.placeholders[1]
                text_frame = content_box.text_frame
                text_frame.clear()
                
                for line in content:
                    line = line.strip()
                    if line:
                        p = text_frame.add_paragraph()
                        p.text = line.lstrip('- *•').strip()
                        p.level = 0 if not line.startswith((' ', '\t')) else 1
            
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                
                # H1 - Title slide
                if line.startswith('# '):
                    if current_slide_title:
                        add_slide(current_slide_title, current_slide_content)
                        current_slide_content = []
                    
                    # Title slide
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
                    
                    title.text = line[2:].strip()
                    if metadata:
                        subtitle.text = f"{metadata.get('company_name', '')}\n{datetime.now().strftime('%B %Y')}"
                    
                    current_slide_title = None
                
                # H2 - New slide
                elif line.startswith('## '):
                    if current_slide_title:
                        add_slide(current_slide_title, current_slide_content)
                        current_slide_content = []
                    
                    current_slide_title = line[3:].strip()
                
                # H3 or content
                elif line.startswith('### '):
                    current_slide_content.append(f"• {line[4:].strip()}")
                
                elif line.startswith(('- ', '* ', '• ')):
                    current_slide_content.append(line)
                
                elif line and current_slide_title:
                    current_slide_content.append(line)
                
                i += 1
            
            # Add last slide
            if current_slide_title:
                add_slide(current_slide_title, current_slide_content)
            
            # Save
            prs.save(output_path)
            logger.info(f"PowerPoint generated: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"PowerPoint conversion failed: {e}")
            raise
    
    def convert_document(self, markdown_content: str, document_id: str, output_dir: str, 
                        metadata: Dict = None, formats: List[str] = None) -> Dict[str, str]:
        """
        Convert a document to multiple formats
        
        Args:
            markdown_content: The markdown content to convert
            document_id: Unique ID for the document
            output_dir: Directory to save files
            metadata: Document metadata (company name, type, etc.)
            formats: List of formats to generate ['pdf', 'word', 'pptx']
        
        Returns:
            Dict of format: file_path
        """
        if formats is None:
            formats = ['pdf', 'word']  # Default formats
        
        os.makedirs(output_dir, exist_ok=True)
        
        results = {}
        base_filename = f"{document_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        try:
            # PDF
            if 'pdf' in formats:
                pdf_path = os.path.join(output_dir, f"{base_filename}.pdf")
                self.convert_to_pdf(markdown_content, pdf_path, metadata)
                results['pdf'] = pdf_path
            
            # Word
            if 'word' in formats or 'docx' in formats:
                word_path = os.path.join(output_dir, f"{base_filename}.docx")
                self.convert_to_word(markdown_content, word_path, metadata)
                results['word'] = word_path
            
            # PowerPoint (primarily for pitch decks)
            if 'pptx' in formats or 'powerpoint' in formats:
                pptx_path = os.path.join(output_dir, f"{base_filename}.pptx")
                self.convert_to_powerpoint(markdown_content, pptx_path, metadata)
                results['pptx'] = pptx_path
            
            logger.info(f"Document converted to {len(results)} formats")
            return results
            
        except Exception as e:
            logger.error(f"Document conversion failed: {e}")
            raise


# Singleton
_converter = None

def get_converter():
    """Get singleton document converter instance"""
    global _converter
    if _converter is None:
        _converter = DocumentConverter()
    return _converter
